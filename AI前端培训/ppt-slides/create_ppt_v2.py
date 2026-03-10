import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# === Colors ===
C_BLUE = RGBColor(0, 102, 255)
C_DARK = RGBColor(30, 41, 59)
C_GRAY = RGBColor(100, 116, 139)
C_BG_LIGHT = RGBColor(248, 250, 252)
C_WHITE = RGBColor(255, 255, 255)
C_ORANGE = RGBColor(245, 158, 11)
C_GREEN = RGBColor(16, 185, 129)
C_RED = RGBColor(239, 68, 68)

def set_font(font, size, color=C_DARK, bold=False):
    font.name = 'Microsoft YaHei'
    font.size = Pt(size)
    font.color.rgb = color
    font.bold = bold

def add_blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_title(slide, text, top=0.5, color=C_BLUE):
    # Indicator bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(top+0.05), Inches(0.12), Inches(0.45))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    
    # Title text
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(10), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    set_font(p.font, 32, C_DARK, True)

def add_text_block(slide, left, top, width, height, text_lines, is_list=True):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text_lines):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = line
        p.level = 1 if is_list else 0
        p.space_after = Pt(12)
        set_font(p.font, 18, C_DARK)

def add_placeholder(slide, left, top, width, height, text="📸 预留配图/截图区"):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_BG_LIGHT
    shape.line.color.rgb = C_GRAY
    shape.line.dash_style = 4
    
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    set_font(p.font, 18, C_GRAY)

def add_card(slide, left, top, width, height, title, lines, title_color=C_BLUE, bg_color=C_WHITE, is_list=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(226, 232, 240)
    
    # top accent line
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = title_color
    accent.line.fill.background()
    
    tf = shape.text_frame
    tf.margin_top = Inches(0.3)
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.3)
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = title
    p.space_after = Pt(10)
    set_font(p.font, 20, title_color, True)
    
    for line in lines:
        p = tf.add_paragraph()
        p.text = line
        p.level = 1 if is_list else 0
        p.space_after = Pt(8)
        set_font(p.font, 16, C_DARK)

def add_chapter_slide(num_text, title, subtitle, color=C_BLUE):
    slide = add_blank_slide()
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(4))
    tf = tb.text_frame
    
    p1 = tf.paragraphs[0]
    p1.text = num_text
    set_font(p1.font, 80, C_WHITE, True)
    
    p2 = tf.add_paragraph()
    p2.text = title
    p2.space_before = Pt(20)
    set_font(p2.font, 48, C_WHITE, True)
    
    p3 = tf.add_paragraph()
    p3.text = subtitle
    p3.space_before = Pt(20)
    set_font(p3.font, 24, C_WHITE)

# ================= PPT Generation =================

# 1. Cover
slide = add_blank_slide()
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = C_BG_LIGHT
bg.line.fill.background()

tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(6.5), Inches(2))
tb.text_frame.word_wrap = True
p = tb.text_frame.paragraphs[0]
p.text = "AI 辅助前端开发\n实战分享"
set_font(p.font, 52, C_BLUE, True)

tb_sub = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(6.5), Inches(1))
tb_sub.text_frame.word_wrap = True
p2 = tb_sub.text_frame.paragraphs[0]
p2.text = "工程落地版 · 从认知打通到工作流重塑"
set_font(p2.font, 22, C_GRAY)

tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(6.5), Inches(1))
tb2.text_frame.word_wrap = True
p3 = tb2.text_frame.paragraphs[0]
p3.text = "主讲人：张兵涛 (前端开发工程师)   |   2026 内部培训"
set_font(p3.font, 18, C_DARK)

slide.shapes.add_picture("/Users/krui/.gemini/antigravity/brain/c396399e-c153-4658-9345-5e0dc2ef0fbe/ppt_cover_bg_1773107972396.png", Inches(7.5), Inches(1.5), width=Inches(5.0))

# 2. Intro Case
slide = add_blank_slide()
add_title(slide, "开场案例：5 分钟感受 AI 能力边界")
add_text_block(slide, 0.5, 1.5, 5.5, 5, [
    "真实的系统搭建流程：",
    "• 输入：🎨 Figma 设计稿结构",
    "• 上下文：📄 ApiFox 接口文档契约",
    "• 输出：💻 高还原度成品页面",
    "",
    "黄金法则：",
    "先确认设计与接口契约，再让 AI 接管繁重代码。AI 是执行高手，方向必须由人定。"
], is_list=False)

slide.shapes.add_picture("/Users/krui/.gemini/antigravity/brain/c396399e-c153-4658-9345-5e0dc2ef0fbe/ppt_case_intro_1773107989052.png", Inches(6.5), Inches(1.5), height=Inches(5))


# 3. TOC
slide = add_blank_slide()
add_title(slide, "大纲概览", color=C_DARK)
add_card(slide, 1, 2.5, 3.5, 3.5, "Part 01\n认知与武器库", [
    "解决“是什么”与“用什么”",
    "• 工具的谱系演化",
    "• 能力边界划定",
    "• 三大核心抓手"
], C_BLUE)
add_card(slide, 4.9, 2.5, 3.5, 3.5, "Part 02\n工作流重塑与实操", [
    "解决“怎么用”",
    "• 告别旧习惯痛点",
    "• 从零建页实战演练",
    "• 500行老组件大重构"
], C_ORANGE)
add_card(slide, 8.8, 2.5, 3.5, 3.5, "Part 03\n避坑指南与沉淀", [
    "防踩雷，保效率",
    "• TOP3 致命翻车操作",
    "• 架构控制权陷阱",
    "• 向团队推广落地资产"
], C_GREEN)

# Part 01
add_chapter_slide("01", "认知与武器库", "解决“是什么”与“用什么”", C_BLUE)

# 1.1 Tools
slide = add_blank_slide()
add_title(slide, "01. 工具的演化与本质")
add_card(slide, 0.5, 1.8, 3.8, 4.5, "阶段一：补全式", [
    "代表：GitHub Copilot",
    "强项：内联语法打字机，极速补写单行或函数",
    "短板：无全局项目感知",
    "👉 就像个打字快的小弟"
], C_GRAY)
add_card(slide, 4.7, 1.8, 3.8, 4.5, "阶段二：会话式", [
    "代表：ChatGPT / 文心",
    "强项：交互解题与查错",
    "短板：上下文极易断裂，需人工大量 C+V (复制粘贴)",
    "👉 就像外包高级技术顾问"
], C_ORANGE)
add_card(slide, 8.9, 1.8, 3.8, 4.5, "阶段三：工程化", [
    "代表：Antigravity",
    "强项：具备全仓记忆图谱，实现 Chat -> Edit -> Commit 端到端闭环",
    "👉 就像具备全局视角的工程管家"
], C_BLUE, RGBColor(235,244,255))

# 1.2 Boundary
slide = add_blank_slide()
add_title(slide, "02. 能力边界划定：何时放权？")
add_card(slide, 0.5, 1.8, 6, 4.5, "✅ 大胆交给 AI（脏活累活）", [
    "• 骨架生成：按图层反出 Vue/React 模板",
    "• 批量平替：API迁移、JS转TS、重命名",
    "• 日志查错：堆栈提取报错根源",
    "• 全能辅笔：写注释、枚举冗长的边界用例"
], C_BLUE)
add_card(slide, 6.8, 1.8, 6, 4.5, "⚠️ 一定靠人工主导（核心决策）", [
    "• 需求把脉：业务算命，AI 听不懂潜规则",
    "• 底层架构：组件体系怎么拆，选什么库",
    "• 代码兜底：Code Review的最终拍板权",
    "",
    "【防沉迷准则】：当 AI 在一个点卡了 15 分钟无产出时，必须立刻人工接管寻找新方案。"
], C_ORANGE, RGBColor(255,247,237))

# 1.3 Core
slide = add_blank_slide()
add_title(slide, "03. 掌控 AI 的三大核心抓手")
add_text_block(slide, 0.5, 1.5, 12, 1, ["掌握它们，AI 就从“野生工具”变成“规范资产”。"], is_list=False)
add_card(slide, 0.5, 2.5, 3.8, 4.5, "1. RAG (知识供给)", [
    "将外部的 ApiFox 契约、需求文档喂给 AI",
    "建立当前业务开发的事实防线，防止幻觉乱编接口。"
], C_BLUE)
add_card(slide, 4.7, 2.5, 3.8, 4.5, "2. Skill (规范指令)", [
    "把团队的口口相传的代码规范（不能Any，命名必须CamalCase）固化成强制系统指令",
    "像护栏一样框死 AI 的输出风格。"
], C_GREEN)
add_card(slide, 8.9, 2.5, 3.8, 4.5, "3. Workflow (标准链路)", [
    "针对复杂高频动作，制定标准的流水线工序：",
    "读文档 -> 写测试纲要 -> 依纲要编写代码 -> 清理验证",
    "一键触发标准开发。 "
], C_ORANGE)

# Part 02
add_chapter_slide("02", "工作流重塑与实操", "解决“怎么用”的痛点", C_ORANGE)

# 2.1 Workflow
slide = add_blank_slide()
add_title(slide, "01. 旧习惯痛点 vs 新 6 步法", color=C_ORANGE)
add_card(slide, 0.5, 1.8, 5.5, 5, "👎 过去痛点：口述需求，边猜边写", [
    "• 理解偏差极大，导致严重的二次返工",
    "• 代码样板全靠手搓，手柄耗时长",
    "• 联调痛苦：接口对不上，全版 console.log 抓瞎"
], C_GRAY, RGBColor(241,245,249))
add_card(slide, 6.5, 1.8, 6.5, 5, "🚀 新型 AI 增强 6 步推进法", [
    "① 人工主导：架构设计与需求理解",
    "② 团队对齐：确定技术方案再动工",
    "③ 颗粒拆解：列出极细的任务列表 (最容易跳过)",
    "④ AI 驱动：一 Task 一交付，即时反馈查验",
    "⑤ 接口联调：集成本地接口定义，精准打通",
    "⑥ 高质审核：专注 Code Review 与可维护性质检"
], C_BLUE, RGBColor(235,244,255))

# 2.2 Case A
slide = add_blank_slide()
add_title(slide, "02. 实战 A：从零建页演示", color=C_ORANGE)
add_text_block(slide, 0.5, 1.5, 5.5, 5, [
    "Step 1：喂入需求大纲",
    "告诉 AI 这是个什么页面，有哪些功能模块。",
    "",
    "Step 2：抽取骨架",
    "基于结构搭建出基础组件树 (Card/Table/Form)。",
    "",
    "Step 3：直连契约注入逻辑",
    "依据接口 Schema 一次性配平拦截器、异常、加载状态。",
    "",
    "爽点：无需等后端写完接口，有一份文档前端即刻完成。 "
], is_list=False)
slide.shapes.add_picture("/Users/krui/.gemini/antigravity/brain/c396399e-c153-4658-9345-5e0dc2ef0fbe/media__1773109502896.png", Inches(6.5), Inches(1.5), width=Inches(6.3))

# 2.3 Case B
slide = add_blank_slide()
add_title(slide, "03. 实战 B：存量巨无霸组件重构", color=C_ORANGE)
add_text_block(slide, 0.5, 1.5, 5.5, 5, [
    "应对 500 行无注释老代码的安全套路：",
    "",
    "1. 全盘拉取分析",
    "扫读全组件，生成现存逻辑图谱。",
    "",
    "2. 勒令生成『切分计划』",
    "由 AI 先出计划，【强人工介入评审】确认方案。",
    "",
    "3. 断点式改造与边界加固",
    "切分 State / View / API，重写一部分验证一部分；用 AI 并行补齐极端的边界拦截校验。"
], is_list=False)
slide.shapes.add_picture("/Users/krui/.gemini/antigravity/brain/c396399e-c153-4658-9345-5e0dc2ef0fbe/media__1773109542628.png", Inches(6.5), Inches(1.5), width=Inches(6.3))

# Part 03
add_chapter_slide("03", "避坑指南与团队沉淀", "防踩雷，从个人到团队", C_GREEN)

# 3.1 Pitfalls
slide = add_blank_slide()
add_title(slide, "01. 致命操作避坑 (Top 3 翻车重灾区)", color=C_GREEN)
add_card(slide, 0.5, 1.5, 4, 4.5, "🔥 需求不清就撒手", [
    "最致命行为！",
    "没有清晰文档就让AI自由发挥，垃圾进必定导致垃圾出，南辕北辙带来毁灭级返工重写。"
], C_RED, RGBColor(254,226,226))
add_card(slide, 4.7, 1.5, 4, 4.5, "🍔 一次指令塞太多要求", [
    "让 AI 同时修 10 个地方。",
    "直接导致多源文件上下文崩塌、幻觉乱发、丢失编辑上下文。必须打细颗粒度派发任务。"
], C_ORANGE, RGBColor(255,247,237))
add_card(slide, 8.9, 1.5, 4, 4.5, "⚠️ 盲信直接提交", [
    "对 AI 编辑内容不做测试和肉眼核对直接 Git commit。",
    "必须坚持人工 Code Review 把守最后防线（特别是潜藏的魔法数字和废弃引入）。"
], C_ORANGE, RGBColor(255,247,237))

# 3.2 Control loss
slide = add_blank_slide()
add_title(slide, "02. 架构与控制权的隐性剥夺", color=C_GREEN)
add_text_block(slide, 0.5, 1.5, 5.5, 5, [
    "警惕 AI 的万金油综合症",
    "它倾向推荐通用或最新潮的技术架构，无视你项目盘根错节的历史负债情况。",
    "",
    "黑盒逻辑的技术债灾难",
    "长期靠 AI 修补核心模块，几个月后宕机，连你自己都看不懂这段逻辑是为了啥。",
    "",
    "安全准则",
    "“哪怕是 AI 生成的，你也必须说得出每一行的存在依据。”"
], is_list=False)

slide.shapes.add_picture("/Users/krui/.gemini/antigravity/brain/c396399e-c153-4658-9345-5e0dc2ef0fbe/ppt_pitfalls_warning_1773108006714.png", Inches(6.5), Inches(1.5), width=Inches(6.3))

# 3.3 Promote
slide = add_blank_slide()
add_title(slide, "03. 从个人破冰到团队资产大跃迁", color=C_GREEN)
add_text_block(slide, 0.5, 1.5, 5.5, 5, [
    "拒绝各自“闭门造车”",
    "好用的规则（Prompt）都是被坑喂出来的。",
    "必须沉淀到项目的 .rules 规范池，形成团队护城河。",
    "",
    "推进落地四段论：",
    "① 个人先通关核心全闭环",
    "② 凝练 2-3 个通用痛点 Workflow",
    "③ 半拉半推引入相关人员尝试",
    "④ 微缩迭代全面推广提振团队坪效"
], is_list=False)

slide.shapes.add_picture("/Users/krui/.gemini/antigravity/brain/c396399e-c153-4658-9345-5e0dc2ef0fbe/media__1773111605293.png", Inches(6.5), Inches(1.5), width=Inches(6.3))

# Thanks
# Thanks
slide = add_blank_slide()
add_title(slide, "Q & A", top=2.5, color=C_BLUE)
add_text_block(slide, 0.8, 4.0, 11, 2, [
    "共识认知边界 · 重塑开发流程 · 沉淀共有资产",
    "谢谢聆听！"
], is_list=False)

output_path = "/Users/krui/Documents/企智数创/AI前端培训/AI前端培训/ppt-slides/AI辅助前端开发_最终确稿.pptx"
prs.save(output_path)
print(f"PPTX Successfully created at: {output_path}")
