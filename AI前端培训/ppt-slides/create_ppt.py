from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
# Set to 16:9 ratio
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette
color_primary = RGBColor(0, 102, 255) # Blue
color_text = RGBColor(26, 31, 54) # Dark gray
color_muted = RGBColor(100, 116, 139) # Light gray

def add_title_slide(title_text, subtitle_text):
    slide_layout = prs.slide_layouts[0] 
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = title_text
    subtitle.text = subtitle_text
    
    # Custom styling
    title.text_frame.paragraphs[0].font.color.rgb = color_primary
    title.text_frame.paragraphs[0].font.bold = True
    subtitle.text_frame.paragraphs[0].font.color.rgb = color_muted

def add_content_slide(title_text, content_list, is_chapter_intro=False):
    layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(layout)
    title = slide.shapes.title
    title.text = title_text
    
    # Title styling
    if is_chapter_intro:
        title.text_frame.paragraphs[0].font.color.rgb = color_primary
        # move to center
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    else:
        title.text_frame.paragraphs[0].font.color.rgb = color_text
    
    title.text_frame.paragraphs[0].font.bold = True
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear() # clear default bullet
    
    for item in content_list:
        p = tf.add_paragraph()
        if isinstance(item, tuple):
             # Indicates sub-bullet
             p.text = item[0]
             p.level = 1
             p.font.size = Pt(20)
             p.font.color.rgb = color_muted
        else:
             p.text = item
             p.level = 0
             p.font.size = Pt(24)
             p.font.color.rgb = color_text
             # add space before
             p.space_before = Pt(14)
            
    # Add a visual placeholder box for screenshots
    if not is_chapter_intro:
        left = Inches(8)
        top = Inches(2.5)
        width = Inches(4.5)
        height = Inches(4)
        shape = slide.shapes.add_shape(
            1, left, top, width, height # MSO_SHAPE.RECTANGLE is 1
        )
        shape.text = "[预留图片/截图添加区]"
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        shape.text_frame.paragraphs[0].font.color.rgb = color_muted
        shape.fill.background()
        shape.line.color.rgb = color_muted

# ================= PPT Content Generation =================

# 1. 封面
add_title_slide(
    "AI 辅助前端开发实战分享", 
    "工程落地版 · 从认知打通到工作流重塑\n\n主讲人：张兵涛 (前端开发工程师)\n2026 · 内部培训"
)

# 2. 开场引语
add_content_slide(
    "开场：5 分钟感受 AI 能力边界",
    [
        "用一段真实的系统搭建作为引子：",
        ("输入：🎨 Figma 设计稿结构图"),
        ("上下文：📄 ApiFox 接口文档契约"),
        ("输出：💻 极高还原度成品页面"),
        "核心思考：AI 是执行高手，方向必须由人定",
        ("先确认设计与接口契约，再让 AI 入场并行开发")
    ]
)

# 3. 目录
add_content_slide(
    "目录大纲",
    [
        "Part 01：认知与武器库",
        ("解决“是什么”与“用什么”"),
        "Part 02：工作流重塑与实操",
        ("解决开发过程中的旧习惯痛点，给出“怎么用”"),
        "Part 03：避坑指南与沉淀",
        ("最高危操作警示与防踩雷，从个人到团队")
    ]
)

# Part 01
add_content_slide("Part 01: 认知与武器库", ["01. 工具演化与本质", "02. 能力边界与适用场景", "03. 三大核心抓手"], is_chapter_intro=True)

add_content_slide(
    "01. 工具的演化与本质",
    [
        "阶段一：补全式 (如 Copilot)",
        ("强项：内联语法打字机。缺点：无全局项目感知。"),
        "阶段二：会话式 (如 ChatGPT)",
        ("强项：强问答找错。缺点：上下文易断裂，需人工大量CV。"),
        "阶段三：工程化闭环 (如 Antigravity)",
        ("强项：具备全仓记忆，实现 Chat -> Edit -> Commit 端到端。")
    ]
)

add_content_slide(
    "02. 能力边界划定：何时放权？",
    [
        "✅ 大胆交给 AI（脏活累活）",
        ("骨架生成、大批量API替换/TS化、复杂堆栈查错、写注释。"),
        "⚠️ 一定要靠人工（核心决策）",
        ("底层架构套路、组件体系切分、需求真伪验证。"),
        "防沉迷准则",
        ("当 AI 磨蹭 15 分钟给不出方向时，必须人工介入叫停。")
    ]
)

add_content_slide(
    "03. 掌控 AI 的三大核心抓手",
    [
        "1. RAG (知识供给)",
        ("将接口与设计文档喂给 AI，建立当前领域事实层。"),
        "2. Skill (指令护栏)",
        ("把团队的代码要求（不能Any，命名法则）变成强制指令。"),
        "3. Workflow (标准链路)",
        ("制定标准的推进工序流程，像剧本一样驱动 AI 执行。")
    ]
)

# Part 02
add_content_slide("Part 02: 工作流重塑与实操", ["01. 旧习惯痛点 vs 新 6 步法", "02. 实战 A：从零建页演示", "03. 实战 B：存量庞大组件重构"], is_chapter_intro=True)

add_content_slide(
    "01. 旧习惯的痛点 vs 新型 6 步流水线",
    [
        "告别口述猜想 -> 改为【方案先行，任务拆解再编码】",
        "AI 辅助下的全新六步：",
        ("① 需求理解与方案（人工）"),
        ("② 技术架构团队对齐（人工）"),
        ("③ 颗粒化 Task 计划罗列（核心转折点）"),
        ("④ AI 驱动编码：一任务一提交"),
        ("⑤ MCP 集成 ApiFox 即时联调"),
        ("⑥ 全面 Code Review（把控质量）")
    ]
)

add_content_slide(
    "02. 实战 A：从零建页演示流程",
    [
        "Step 1：喂入需求大纲",
        ("引入 Figma 结构让 AI 理清布局关系"),
        "Step 2：抽取骨架",
        ("直出全套空白 UI 组件树代码"),
        "Step 3：直连契约生成处理",
        ("引入 Swagger/API 结构，自动配平字段处理方案"),
        "结论体验",
        ("前端无需苦等后端进度接口布完才可联调")
    ]
)

add_content_slide(
    "03. 实战 B：面对 500 行的遗留代码怎么改",
    [
        "历史欠债的安全切块",
        ("先让 AI 生成整体依赖思维导图，理清脉络"),
        "强人工干预",
        ("AI 给出《拆借计划》-> 人工确认/否决 -> 再开始生成代码"),
        "断点加固测试",
        ("用 AI 一并补齐历史空白的边界异常数据处理（如网络错误空降落伞）")
    ]
)

# Part 03
add_content_slide("Part 03: 避坑指南与经验沉淀", ["01. 致命操作的防范", "02. 架构控制权的陷阱", "03. 如何实现向团队推广"], is_chapter_intro=True)

add_content_slide(
    "01. 致命操作避坑 (Top 3 翻车重灾区)",
    [
        "坑 1：需求连自己都没理顺",
        ("垃圾需求输入 -> 高难度乱涂输出 -> 巨额维护灾难"),
        "坑 2：指派的活儿包太厚",
        ("一次给 AI 下达十几个分散的修改命令极易造成上下文崩塌幻觉"),
        "坑 3：盲信提交，免除把关",
        ("直接把 AI 的修改点推上线，缺乏 Review 的黑盒极具毁灭性")
    ]
)

add_content_slide(
    "02. 警惕：工程架构控制权剥夺",
    [
        "AI 的万金油陷阱",
        ("它随时可能抛出并不符合本项目旧基建的时髦重构方案"),
        "黑盒逻辑的积累",
        ("无人懂这段连线代码是干嘛的，3个月后就是个无法维护的雷"),
        "解决准则",
        ("必须始终“能说明白 AI 写的每行判断依据”，否则不予通过")
    ]
)

add_content_slide(
    "03. 从个人破冰到团队推广沉淀",
    [
        "拒绝一个人造车",
        ("个人的零碎 Prompt 应该汇总在 .rules 或公库资产中"),
        "标准化复用",
        ("提炼团队最高频痛点，建立 2-3 套公用 Workflow 护航"),
        "落地推进节奏",
        ("个人通关 -> 分享小圈 -> 半年微迭代推广全组")
    ]
)

# 最后谢幕
add_title_slide("Q & A", "共识认知边界 · 重塑开发流程 · 沉淀共有资产\n\n谢谢大家！")

output_path = "/Users/krui/Documents/企智数创/AI前端培训/AI前端培训/ppt-slides/AI辅助前端开发_张兵涛_分享.pptx"
prs.save(output_path)
print(f"PPTX Successfully created at: {output_path}")
